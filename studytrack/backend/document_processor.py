import io
import re
from collections import Counter

import fitz
from pptx import Presentation


def _normalize_text(value: str) -> str:
    # Collapse noisy whitespace into predictable single spaces.
    return re.sub(r"\s+", " ", value or "").strip()


def extract_text_from_pdf(file_bytes: bytes) -> str:
    document = fitz.open(stream=file_bytes, filetype="pdf")
    all_text: list[str] = []

    for page in document:
        page_text = _normalize_text(page.get_text("text"))
        if page_text:
            all_text.append(page_text)

    return "\n\n".join(all_text)


def extract_text_from_pptx(file_bytes: bytes) -> str:
    presentation = Presentation(io.BytesIO(file_bytes))
    slides_text: list[str] = []

    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = [f"--- Slide {index} ---"]

        for shape in slide.shapes:
            text = None
            if getattr(shape, "has_text_frame", False) and shape.text_frame:
                text = shape.text_frame.text
            elif hasattr(shape, "text"):
                text = shape.text

            cleaned = _normalize_text(text or "")
            if cleaned:
                parts.append(cleaned)

        slides_text.append("\n".join(parts))

    return "\n\n".join(slides_text)


def get_relevant_chunks(topic_name: str, all_chunks: list[str], max_chunks: int = 5) -> list[str]:
    if not all_chunks:
        return []

    query_tokens = [
        token.lower()
        for token in re.findall(r"[A-Za-z0-9]+", topic_name or "")
        if len(token) > 1
    ]

    if not query_tokens:
        return all_chunks[:max_chunks]

    token_counter = Counter(query_tokens)
    scored: list[tuple[int, str]] = []

    for chunk in all_chunks:
        haystack = (chunk or "").lower()
        score = sum(haystack.count(token) * weight for token, weight in token_counter.items())
        scored.append((score, chunk))

    scored.sort(key=lambda item: item[0], reverse=True)
    top = [chunk for score, chunk in scored if score > 0][:max_chunks]

    if not top:
        return all_chunks[:max_chunks]
    return top
